import io

from app.models.enums import SourceType
from app.parsing.base import ParsedDocument, TextBlock, register


@register(SourceType.ppt)
class PptxParser:
    def parse(self, *, data: bytes | None, url: str | None, filename: str) -> ParsedDocument:
        from pptx import Presentation

        if data is None:
            raise ValueError("PPTX parser requires file bytes")
        prs = Presentation(io.BytesIO(data))
        blocks: list[TextBlock] = []
        for i, slide in enumerate(prs.slides):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        run_text = "".join(run.text for run in para.runs)
                        if run_text.strip():
                            parts.append(run_text)
            if parts:
                blocks.append(TextBlock(text="\n".join(parts), page_number=i + 1))
        return ParsedDocument(
            blocks=blocks, page_count=len(prs.slides), meta={"parser": "python-pptx"}
        )
